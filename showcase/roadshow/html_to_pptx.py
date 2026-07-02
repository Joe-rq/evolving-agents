# -*- coding: utf-8 -*-
"""
痕进路演 HTML → PPTX 转换器
--------------------------
把 roadshow/index.html 的 10 个 <section class="slide"> 逐页用本机 Chrome
headless 截图（100% 保留现有 CSS 设计），再用 python-pptx 拼成 16:9 .pptx。

零新增依赖：只用系统已装的 Chrome + python-pptx。
"""
import os
import re
import shutil
import subprocess
import sys
import tempfile

from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
SRC_HTML = os.path.join(HERE, "index.html")
OUT_PPTX = os.path.join(HERE, "痕进_路演.pptx")

CHROME = r"C:\Program Files\Google\Chrome\Application\chrome.exe"
CHROME_X86 = r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"

# 渲染分辨率：1920x1080 (16:9)，scale=2 → 输出 3840x2160 高清
VIEW_W, VIEW_H, SCALE = 1920, 1080, 2


def find_chrome():
    for p in (CHROME, CHROME_X86):
        if os.path.exists(p):
            return p
    # 兜底：PATH 里的 chrome
    return shutil.which("chrome") or shutil.which("google-chrome")


def extract(src_path):
    """从 index.html 抽取 <head>(含 <style>) 与每个 <section class="slide">。"""
    html = open(src_path, encoding="utf-8").read()
    head = re.search(r"<head>.*?</head>", html, re.S).group(0)
    # 所有幻灯片段
    sections = re.findall(
        r'<section class="slide"[^>]*>.*?</section>', html, re.S)
    if not sections:
        sys.exit("未找到 <section class=\"slide\"> 片段")
    return head, sections


def render_slide(chrome, head, section, idx, out_png):
    """为单个 section 生成独立 HTML（去掉 deck-nav，单页满屏），Chrome 截图。"""
    # 把 #x 的 id 改成 s1，确保 :target / 满屏逻辑不变；section 已自带 id
    standalone = (
        "<!DOCTYPE html>\n<html lang=\"zh-CN\">\n"
        + head
        + "\n<body>\n"
        + section
        + "\n</body>\n</html>\n"
    )
    tmp_html = os.path.join(HERE, "build", f"slide_{idx:02d}.html")
    with open(tmp_html, "w", encoding="utf-8") as f:
        f.write(standalone)

    file_url = "file:///" + tmp_html.replace("\\", "/")
    cmd = [
        chrome,
        "--headless=new",
        "--disable-gpu",
        "--no-sandbox",
        "--hide-scrollbars",
        "--force-device-scale-factor=%d" % SCALE,
        "--window-size=%d,%d" % (VIEW_W, VIEW_H),
        "--default-background-color=00000000",
        "--screenshot=%s" % out_png,
        file_url,
    ]
    r = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
    if not os.path.exists(out_png):
        sys.exit("截图失败 slide %d:\n%s\n%s" % (idx, r.stdout, r.stderr))


def build_pptx(pngs):
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            png, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(OUT_PPTX)


def main():
    chrome = find_chrome()
    if not chrome:
        sys.exit("找不到 Chrome，请安装或调整脚本里的 CHROME 路径")
    print("Chrome :", chrome)

    build_dir = os.path.join(HERE, "build")
    if os.path.isdir(build_dir):
        shutil.rmtree(build_dir)
    os.makedirs(build_dir, exist_ok=True)

    head, sections = extract(SRC_HTML)
    print("幻灯片数:", len(sections))

    pngs = []
    for i, sec in enumerate(sections, 1):
        out_png = os.path.join(HERE, "build", f"slide_{i:02d}.png")
        render_slide(chrome, head, sec, i, out_png)
        pngs.append(out_png)
        print("  [%02d] 截图完成 → %s" % (i, os.path.basename(out_png)))

    build_pptx(pngs)
    print("\n✅ 已生成:", OUT_PPTX)
    print("   共 %d 页，16:9 (3840x2160 高清源图)" % len(pngs))


if __name__ == "__main__":
    main()
